"""
テナントB Functions → テナントA Graph (SharePoint) アクセス
認証方式: Managed Identity + Workload Identity Federation

認証フロー:
  1. UAMI (User Assigned Managed Identity) で api://AzureADTokenExchange トークンを取得
  2. そのトークンを Client Assertion として テナントA 向け Graph トークンを取得
  3. Graph API で テナントA の SharePoint サイトへアクセス
"""

import hashlib
import io
import json
import logging
import os
import re
import time
from pathlib import Path
from typing import Any, Callable, Dict, List, Optional, Tuple

import azure.functions as func
import requests
from azure.identity import ClientAssertionCredential, ManagedIdentityCredential

# ドキュメント処理ライブラリ (オプション)
try:
    from pdfminer.high_level import extract_text as pdf_extract_text
    from pdfminer.layout import LAParams
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from openpyxl import load_workbook
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False

app = func.FunctionApp(http_auth_level=func.AuthLevel.FUNCTION)
GRAPH_BASE_URL = "https://graph.microsoft.com/v1.0"
GRAPH_SCOPE = "https://graph.microsoft.com/.default"
FEDERATION_SCOPE = "api://AzureADTokenExchange/.default"


def _get_env(name: str, required: bool = True, default: Optional[str] = None) -> str:
    value = os.getenv(name, default)
    if required and (value is None or value.strip() == ""):
        raise ValueError(f"Missing required environment variable: {name}")
    return (value or "").strip()


def _build_assertion_func(mi_client_id: Optional[str] = None) -> Callable[[], str]:
    """Managed Identity で api://AzureADTokenExchange トークンを取得する関数を返す"""
    if mi_client_id:
        mi_credential = ManagedIdentityCredential(client_id=mi_client_id)
    else:
        mi_credential = ManagedIdentityCredential()

    def _get_assertion() -> str:
        token = mi_credential.get_token(FEDERATION_SCOPE)
        return token.token

    return _get_assertion


def _get_graph_token() -> str:
    """
    Workload Identity Federation でクロステナント Graph トークンを取得する。
    - TENANT_A_ID         : アクセス先 (テナントA) のテナント ID
    - APP_CLIENT_ID       : テナントB で作成した Multitenant アプリの Client ID
    - MANAGED_IDENTITY_CLIENT_ID : テナントB の UAMI Client ID (省略時はシステム割当)
    """
    tenant_a_id = _get_env("TENANT_A_ID")
    app_client_id = _get_env("APP_CLIENT_ID")
    mi_client_id = os.getenv("MANAGED_IDENTITY_CLIENT_ID", "").strip() or None

    assertion_func = _build_assertion_func(mi_client_id)

    credential = ClientAssertionCredential(
        tenant_id=tenant_a_id,
        client_id=app_client_id,
        func=assertion_func,
    )
    token = credential.get_token(GRAPH_SCOPE)
    return token.token


def _graph_get(
    path: str, token: str, params: Optional[Dict[str, Any]] = None
) -> Dict[str, Any]:
    headers = {"Authorization": f"Bearer {token}"}
    url = f"{GRAPH_BASE_URL}{path}"

    for attempt in range(4):
        response = requests.get(url, headers=headers, params=params, timeout=30)
        if response.status_code in (429, 503, 504) and attempt < 3:
            retry_after = int(response.headers.get("Retry-After", "2"))
            time.sleep(min(10, retry_after + attempt))
            continue
        if response.status_code >= 400:
            raise RuntimeError(
                f"Graph request failed: {response.status_code} {response.text}"
            )
        return response.json()

    raise RuntimeError("Graph request failed after retries")


def _resolve_site(hostname: str, site_path: str, token: str) -> Dict[str, Any]:
    return _graph_get(
        f"/sites/{hostname}:{site_path}",
        token,
        params={"$select": "id,displayName,webUrl"},
    )


def _list_root_children(site_id: str, token: str, top: int) -> List[Dict[str, Any]]:
    data = _graph_get(
        f"/sites/{site_id}/drive/root/children",
        token,
        params={
            "$top": str(top),
            "$select": "id,name,webUrl,size,createdDateTime,lastModifiedDateTime,eTag,cTag,file,folder,parentReference",
        },
    )
    return data.get("value", [])


def _scan_sites() -> Dict[str, Any]:
    hostname = _get_env("SP_HOSTNAME")
    site_paths_raw = _get_env("SP_SITE_PATHS")
    site_paths = [p.strip() for p in site_paths_raw.split(",") if p.strip()]
    if not site_paths:
        raise ValueError("SP_SITE_PATHS must include at least one site path")

    max_items = int(os.getenv("MAX_ITEMS_PER_SITE", "20"))
    token = _get_graph_token()

    sites_result: List[Dict[str, Any]] = []
    for site_path in site_paths:
        site = _resolve_site(hostname, site_path, token)
        children = _list_root_children(site["id"], token, max_items)
        file_items = []

        for item in children:
            file_items.append(
                {
                    "id": item.get("id"),
                    "title": item.get("name"),
                    "url": item.get("webUrl"),
                    "size": item.get("size"),
                    "createdDateTime": item.get("createdDateTime"),
                    "lastModifiedDateTime": item.get("lastModifiedDateTime"),
                    "isFolder": "folder" in item,
                    "parentPath": item.get("parentReference", {}).get("path"),
                    "eTag": item.get("eTag"),
                }
            )

        sites_result.append(
            {
                "sitePath": site_path,
                "siteId": site.get("id"),
                "siteTitle": site.get("displayName"),
                "siteUrl": site.get("webUrl"),
                "items": file_items,
            }
        )

    return {
        "authMode": "workload_identity_federation",
        "tenantA": _get_env("TENANT_A_ID"),
        "hostname": hostname,
        "siteCount": len(sites_result),
        "sites": sites_result,
    }


# ---------------------------------------------------------------------------
# テキスト抽出
# ---------------------------------------------------------------------------

class TextExtractor:
    """ドキュメントからテキストを抽出"""

    SUPPORTED_EXTENSIONS = [".pdf", ".docx", ".pptx", ".ppt", ".xlsx", ".xls", ".txt"]

    @staticmethod
    def extract(content: bytes, file_extension: str) -> str:
        ext = file_extension.lower()
        if ext == ".pdf":
            return TextExtractor._pdf(content)
        elif ext == ".docx":
            return TextExtractor._docx(content)
        elif ext in (".pptx", ".ppt"):
            return TextExtractor._pptx(content)
        elif ext in (".xlsx", ".xls"):
            return TextExtractor._xlsx(content)
        elif ext == ".txt":
            return TextExtractor._txt(content)
        else:
            raise ValueError(f"Unsupported extension: {ext}")

    @staticmethod
    def _pdf(content: bytes) -> str:
        if not PDF_AVAILABLE:
            raise ImportError("pdfminer.six is required")
        return pdf_extract_text(io.BytesIO(content), laparams=LAParams()).strip()

    @staticmethod
    def _docx(content: bytes) -> str:
        if not DOCX_AVAILABLE:
            raise ImportError("python-docx is required")
        doc = DocxDocument(io.BytesIO(content))
        return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())

    @staticmethod
    def _pptx(content: bytes) -> str:
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx is required")
        prs = Presentation(io.BytesIO(content))
        parts: List[str] = []
        for i, slide in enumerate(prs.slides, 1):
            lines = [f"--- Slide {i} ---"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    lines.append(shape.text.strip())
            parts.append("\n".join(lines))
        return "\n\n".join(parts)

    @staticmethod
    def _xlsx(content: bytes) -> str:
        if not EXCEL_AVAILABLE:
            raise ImportError("openpyxl is required")
        wb = load_workbook(io.BytesIO(content), data_only=True)
        parts: List[str] = []
        for name in wb.sheetnames:
            sheet = wb[name]
            lines = [f"--- Sheet: {name} ---"]
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join(str(c) if c is not None else "" for c in row)
                if row_text.strip(" |"):
                    lines.append(row_text)
            if len(lines) > 1:
                parts.append("\n".join(lines))
        return "\n\n".join(parts)

    @staticmethod
    def _txt(content: bytes) -> str:
        for enc in ("utf-8", "shift-jis", "cp932"):
            try:
                return content.decode(enc)
            except UnicodeDecodeError:
                continue
        return content.decode("utf-8", errors="ignore")


# ---------------------------------------------------------------------------
# テキストチャンク分割
# ---------------------------------------------------------------------------

def _split_text_into_chunks(
    text: str, chunk_size: int = 1000, chunk_overlap: int = 200
) -> List[Dict[str, Any]]:
    """テキストを段落→文の順でチャンクに分割し、リストで返す。"""
    if not text or not text.strip():
        return []

    paragraphs = [p.strip() for p in re.split(r"\n\s*\n", text) if p.strip()]

    chunks: List[Dict[str, Any]] = []
    current = ""
    idx = 0

    def _flush(cur: str) -> str:
        nonlocal idx
        chunks.append({"chunk_index": idx, "text": cur.strip()})
        idx += 1
        if len(cur) > chunk_overlap:
            return cur[-chunk_overlap:]
        return cur

    for para in paragraphs:
        if len(para) > chunk_size:
            sentences = re.split(r"[。.!?]\s*", para)
            for sent in sentences:
                sent = sent.strip()
                if not sent:
                    continue
                if len(current) + len(sent) <= chunk_size:
                    current += sent + " "
                else:
                    if current:
                        current = _flush(current)
                    current += sent + " "
        else:
            if len(current) + len(para) <= chunk_size:
                current += para + "\n\n"
            else:
                if current:
                    current = _flush(current)
                current += para + "\n\n"

    if current.strip():
        chunks.append({"chunk_index": idx, "text": current.strip()})

    return chunks


# ---------------------------------------------------------------------------
# ACL 抽出
# ---------------------------------------------------------------------------

def _extract_acl_from_permissions(
    permissions: List[Dict[str, Any]],
) -> Tuple[List[str], List[str]]:
    """Graph API permissions レスポンスから (user_ids, group_ids) を返す。"""
    user_ids: List[str] = []
    group_ids: List[str] = []

    for perm in permissions:
        granted = perm.get("grantedToV2") or perm.get("grantedTo")
        if granted:
            user = granted.get("user")
            if user and "id" in user:
                user_ids.append(user["id"])
            group = granted.get("group")
            if group and "id" in group:
                group_ids.append(group["id"])

        identities = perm.get("grantedToIdentitiesV2") or perm.get(
            "grantedToIdentities", []
        )
        for identity in identities:
            user = identity.get("user")
            if user and "id" in user:
                user_ids.append(user["id"])
            group = identity.get("group")
            if group and "id" in group:
                group_ids.append(group["id"])

    return list(set(user_ids)), list(set(group_ids))


# ---------------------------------------------------------------------------
# ドキュメント ID 生成
# ---------------------------------------------------------------------------

def _create_document_id(
    site_id: str, drive_id: str, item_id: str, chunk_index: Optional[int] = None
) -> str:
    def _sanitize(text: str, max_len: int = 30) -> str:
        safe = re.sub(r"[^a-zA-Z0-9\-_]", "", text)
        if safe and not safe[0].isalpha():
            safe = "doc" + safe
        if not safe:
            safe = "unknown"
        return safe[:max_len].rstrip("-_") or "unknown"

    s, d, i = _sanitize(site_id, 20), _sanitize(drive_id, 20), _sanitize(item_id, 40)
    doc_id = f"{s}_{d}_{i}_{chunk_index}" if chunk_index is not None else f"{s}_{d}_{i}"
    if len(doc_id) > 1024:
        h = hashlib.md5(doc_id.encode()).hexdigest()[:8]
        doc_id = (
            f"{s[:10]}_{d[:10]}_{h}_{chunk_index}"
            if chunk_index is not None
            else f"{s[:10]}_{d[:10]}_{h}"
        )
    return doc_id


# ---------------------------------------------------------------------------
# ドキュメント取り込み処理 (AI Search アップロードなし)
# ---------------------------------------------------------------------------

def _get_all_files_recursive(
    token: str,
    drive_id: str,
    folder_path: str = "root",
    supported_extensions: Optional[List[str]] = None,
    max_files: int = 50,
    max_depth: int = 5,
    current_depth: int = 0,
    collected: Optional[List[Dict[str, Any]]] = None,
) -> List[Dict[str, Any]]:
    """ドライブ内のファイルを再帰的に取得"""
    if collected is None:
        collected = []
    if current_depth >= max_depth or len(collected) >= max_files:
        return collected

    try:
        data = _graph_get(
            f"/drives/{drive_id}/{folder_path}/children",
            token,
            params={
                "$top": "200",
                "$select": "id,name,size,webUrl,file,folder,parentReference,"
                "lastModifiedDateTime,createdBy,lastModifiedBy",
            },
        )
        for item in data.get("value", []):
            if len(collected) >= max_files:
                break
            name = item["name"]
            if "file" in item:
                ext = Path(name).suffix.lower()
                if supported_extensions is None or ext in supported_extensions:
                    collected.append(
                        {
                            "id": item["id"],
                            "name": name,
                            "path": item.get("parentReference", {}).get("path", "")
                            + "/"
                            + name,
                            "extension": ext,
                            "size": item.get("size", 0),
                            "webUrl": item.get("webUrl", ""),
                            "lastModified": item.get("lastModifiedDateTime", ""),
                            "createdBy": item.get("createdBy", {})
                            .get("user", {})
                            .get("displayName", ""),
                            "modifiedBy": item.get("lastModifiedBy", {})
                            .get("user", {})
                            .get("displayName", ""),
                            "drive_id": drive_id,
                        }
                    )
            elif "folder" in item:
                _get_all_files_recursive(
                    token,
                    drive_id,
                    f"items/{item['id']}",
                    supported_extensions,
                    max_files,
                    max_depth,
                    current_depth + 1,
                    collected,
                )
    except Exception as e:
        logging.warning("Folder listing error (%s): %s", folder_path, str(e)[:200])

    return collected


def _download_file_content(token: str, drive_id: str, item_id: str) -> bytes:
    """Graph API 経由でファイルをダウンロード"""
    url = f"{GRAPH_BASE_URL}/drives/{drive_id}/items/{item_id}/content"
    headers = {"Authorization": f"Bearer {token}"}
    resp = requests.get(url, headers=headers, allow_redirects=True, timeout=60)
    resp.raise_for_status()
    return resp.content


def _get_item_permissions(
    token: str, drive_id: str, item_id: str
) -> Tuple[List[str], List[str]]:
    """ファイルの ACL (permissions) を取得して (user_ids, group_ids) を返す"""
    try:
        data = _graph_get(
            f"/drives/{drive_id}/items/{item_id}/permissions", token
        )
        return _extract_acl_from_permissions(data.get("value", []))
    except Exception as e:
        logging.warning("Permission fetch error (item %s): %s", item_id, str(e)[:200])
        return [], []


def _ingest_sites(
    max_files_per_site: int = 20,
    chunk_size: int = 1000,
    chunk_overlap: int = 200,
    max_depth: int = 5,
    content_preview_len: int = 80,
) -> Dict[str, Any]:
    """
    サイトごとにファイルを取得→テキスト抽出→チャンク分割→ACL取得し、
    ドキュメントオブジェクトの一覧を返す (AI Search へのアップロードは行わない)。
    """
    hostname = _get_env("SP_HOSTNAME")
    site_paths_raw = _get_env("SP_SITE_PATHS")
    site_paths = [p.strip() for p in site_paths_raw.split(",") if p.strip()]
    if not site_paths:
        raise ValueError("SP_SITE_PATHS must include at least one site path")

    token = _get_graph_token()

    sites_result: List[Dict[str, Any]] = []
    total_files = 0
    total_documents = 0
    total_errors = 0

    for site_path in site_paths:
        site = _resolve_site(hostname, site_path, token)
        site_id = site["id"]
        site_name = site.get("displayName", site_path)

        # サイトのドライブ一覧を取得
        drives_data = _graph_get(f"/sites/{site_id}/drives", token)
        drives = drives_data.get("value", [])

        site_docs: List[Dict[str, Any]] = []
        site_errors: List[Dict[str, Any]] = []

        for drive in drives:
            drive_id = drive["id"]
            drive_name = drive.get("name", "Unknown")

            # ファイル一覧を再帰取得
            files = _get_all_files_recursive(
                token,
                drive_id,
                "root",
                supported_extensions=TextExtractor.SUPPORTED_EXTENSIONS,
                max_files=max_files_per_site,
                max_depth=max_depth,
            )

            for file_info in files:
                total_files += 1
                try:
                    # ファイルダウンロード
                    content_bytes = _download_file_content(
                        token, drive_id, file_info["id"]
                    )

                    # テキスト抽出
                    text = TextExtractor.extract(
                        content_bytes, file_info["extension"]
                    )
                    if not text or len(text.strip()) < 10:
                        site_errors.append(
                            {
                                "file": file_info["name"],
                                "reason": "empty or too short text",
                            }
                        )
                        total_errors += 1
                        continue

                    # ACL 取得
                    acl_users, acl_groups = _get_item_permissions(
                        token, drive_id, file_info["id"]
                    )

                    # チャンク分割
                    chunks = _split_text_into_chunks(
                        text, chunk_size, chunk_overlap
                    )

                    for chunk in chunks:
                        doc_id = _create_document_id(
                            site_id, drive_id, file_info["id"], chunk["chunk_index"]
                        )
                        doc = {
                            "id": doc_id,
                            "title": file_info["name"],
                            "contentPreview": chunk["text"][:content_preview_len],
                            "contentLength": len(chunk["text"]),
                            "url": file_info["webUrl"],
                            "path": file_info["path"],
                            "site": site_name,
                            "library": drive_name,
                            "contentType": file_info["extension"][1:]
                            if file_info["extension"]
                            else "unknown",
                            "fileExtension": file_info["extension"],
                            "lastModified": file_info["lastModified"],
                            "createdBy": file_info["createdBy"],
                            "modifiedBy": file_info["modifiedBy"],
                            "size": file_info["size"],
                            "chunkIndex": chunk["chunk_index"],
                            "aclUsers": acl_users,
                            "aclGroups": acl_groups,
                        }
                        site_docs.append(doc)
                        total_documents += 1

                except Exception as e:
                    total_errors += 1
                    site_errors.append(
                        {"file": file_info["name"], "reason": str(e)[:200]}
                    )
                    logging.warning(
                        "File processing error (%s): %s",
                        file_info["name"],
                        str(e)[:200],
                    )

        sites_result.append(
            {
                "sitePath": site_path,
                "siteId": site_id,
                "siteTitle": site_name,
                "siteUrl": site.get("webUrl"),
                "driveCount": len(drives),
                "documents": site_docs,
                "errors": site_errors,
            }
        )

    return {
        "authMode": "workload_identity_federation",
        "tenantA": _get_env("TENANT_A_ID"),
        "hostname": hostname,
        "settings": {
            "maxFilesPerSite": max_files_per_site,
            "chunkSize": chunk_size,
            "chunkOverlap": chunk_overlap,
            "maxDepth": max_depth,
            "contentPreviewLen": content_preview_len,
        },
        "summary": {
            "siteCount": len(sites_result),
            "totalFiles": total_files,
            "totalDocuments": total_documents,
            "totalErrors": total_errors,
        },
        "sites": sites_result,
    }


# ===========================================================================
# HTTP / Timer Triggers
# ===========================================================================


@app.function_name(name="GraphCrossTenantIngestHttp")
@app.route(route="graph/cross-tenant-ingest", methods=["GET"])
def graph_cross_tenant_ingest_http(req: func.HttpRequest) -> func.HttpResponse:
    """
    HTTP Trigger: テナントA の SharePoint サイトからドキュメントを取得・処理し、
    チャンク化されたドキュメント一覧を返す (AI Search アップロードなし)。

    クエリパラメータ:
      max_files   - サイトあたりの最大ファイル数 (default: 10)
      chunk_size  - チャンクサイズ文字数 (default: 1000)
      overlap     - チャンクオーバーラップ文字数 (default: 200)
      max_depth   - フォルダ探索の最大階層 (default: 5)
      preview_len - content プレビューの文字数 (default: 80)
    """
    try:
        max_files = int(req.params.get("max_files", "10"))
        chunk_size = int(req.params.get("chunk_size", "1000"))
        overlap = int(req.params.get("overlap", "200"))
        max_depth = int(req.params.get("max_depth", "5"))
        preview_len = int(req.params.get("preview_len", "80"))

        result = _ingest_sites(
            max_files_per_site=max_files,
            chunk_size=chunk_size,
            chunk_overlap=overlap,
            max_depth=max_depth,
            content_preview_len=preview_len,
        )
        return func.HttpResponse(
            body=json.dumps(result, ensure_ascii=False, indent=2),
            status_code=200,
            mimetype="application/json",
        )
    except Exception as exc:
        logging.exception("GraphCrossTenantIngestHttp failed")
        return func.HttpResponse(
            body=json.dumps({"error": str(exc)}, ensure_ascii=False),
            status_code=500,
            mimetype="application/json",
        )


@app.function_name(name="GraphCrossTenantScanHttp")
@app.route(route="graph/cross-tenant-scan", methods=["GET"])
def graph_cross_tenant_scan_http(req: func.HttpRequest) -> func.HttpResponse:
    """HTTP Trigger: テナントA の SharePoint サイトをスキャンして結果を返す"""
    try:
        result = _scan_sites()
        return func.HttpResponse(
            body=json.dumps(result, ensure_ascii=False, indent=2),
            status_code=200,
            mimetype="application/json",
        )
    except Exception as exc:
        logging.exception("GraphCrossTenantScanHttp failed")
        return func.HttpResponse(
            body=json.dumps({"error": str(exc)}, ensure_ascii=False),
            status_code=500,
            mimetype="application/json",
        )


@app.function_name(name="GraphCrossTenantScanTimer")
@app.schedule(
    schedule="%GRAPH_SCAN_SCHEDULE%",
    arg_name="timer",
    run_on_startup=False,
    use_monitor=True,
)
def graph_cross_tenant_scan_timer(timer: func.TimerRequest) -> None:
    """Timer Trigger: 定期実行で テナントA の SharePoint サイトをスキャン"""
    try:
        result = _scan_sites()
        logging.info(
            "GraphCrossTenantScanTimer completed. siteCount=%s tenantA=%s",
            result.get("siteCount"),
            result.get("tenantA"),
        )
        for site in result.get("sites", []):
            logging.info(
                "site=%s items=%s",
                site.get("sitePath"),
                len(site.get("items", [])),
            )
    except Exception:
        logging.exception("GraphCrossTenantScanTimer failed")
        raise
