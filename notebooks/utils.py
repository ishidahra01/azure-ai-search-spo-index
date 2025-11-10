"""
SharePoint から Azure AI Search への取り込みユーティリティ

このモジュールには、以下の機能が含まれます:
- テキスト抽出 (PDF, Word, PowerPoint, Excel)
- チャンク分割
- 埋め込み生成
- ACL 処理
"""

import io
import re
from pathlib import Path
from typing import List, Dict, Optional, Tuple
from dataclasses import dataclass

# PDF処理
try:
    from pdfminer.high_level import extract_text as pdf_extract_text
    from pdfminer.layout import LAParams
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

# Word処理
try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

# PowerPoint処理
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# Excel処理
try:
    from openpyxl import load_workbook
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False


@dataclass
class TextChunk:
    """テキストチャンク"""
    text: str
    chunk_index: int
    start_char: int
    end_char: int
    metadata: Dict = None


class TextExtractor:
    """ドキュメントからテキストを抽出"""
    
    @staticmethod
    def extract_from_pdf(content: bytes) -> str:
        """PDFからテキスト抽出"""
        if not PDF_AVAILABLE:
            raise ImportError("pdfminer.six が必要です")
        
        try:
            pdf_file = io.BytesIO(content)
            laparams = LAParams()
            text = pdf_extract_text(pdf_file, laparams=laparams)
            return text.strip()
        except Exception as e:
            raise Exception(f"PDF抽出エラー: {str(e)}")
    
    @staticmethod
    def extract_from_docx(content: bytes) -> str:
        """Wordからテキスト抽出"""
        if not DOCX_AVAILABLE:
            raise ImportError("python-docx が必要です")
        
        try:
            doc = DocxDocument(io.BytesIO(content))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n\n".join(paragraphs)
        except Exception as e:
            raise Exception(f"Word抽出エラー: {str(e)}")
    
    @staticmethod
    def extract_from_pptx(content: bytes) -> str:
        """PowerPointからテキスト抽出"""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx が必要です")
        
        try:
            prs = Presentation(io.BytesIO(content))
            text_runs = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"--- スライド {slide_num} ---"]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                text_runs.append("\n".join(slide_text))
            
            return "\n\n".join(text_runs)
        except Exception as e:
            raise Exception(f"PowerPoint抽出エラー: {str(e)}")
    
    @staticmethod
    def extract_from_xlsx(content: bytes) -> str:
        """Excelからテキスト抽出"""
        if not EXCEL_AVAILABLE:
            raise ImportError("openpyxl が必要です")
        
        try:
            wb = load_workbook(io.BytesIO(content), data_only=True)
            text_parts = []
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheet_text = [f"--- シート: {sheet_name} ---"]
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip(" |"):
                        sheet_text.append(row_text)
                
                if len(sheet_text) > 1:
                    text_parts.append("\n".join(sheet_text))
            
            return "\n\n".join(text_parts)
        except Exception as e:
            raise Exception(f"Excel抽出エラー: {str(e)}")
    
    @staticmethod
    def extract_from_txt(content: bytes) -> str:
        """テキストファイルから抽出"""
        try:
            # UTF-8を試す
            return content.decode('utf-8')
        except UnicodeDecodeError:
            try:
                # Shift-JISを試す
                return content.decode('shift-jis')
            except UnicodeDecodeError:
                # CP932を試す
                return content.decode('cp932', errors='ignore')
    
    @staticmethod
    def extract(content: bytes, file_extension: str) -> str:
        """
        ファイル拡張子に基づいてテキスト抽出
        
        Args:
            content: ファイルのバイト内容
            file_extension: 拡張子 (.pdf, .docx など)
        
        Returns:
            抽出されたテキスト
        """
        ext = file_extension.lower()
        
        if ext == '.pdf':
            return TextExtractor.extract_from_pdf(content)
        elif ext == '.docx':
            return TextExtractor.extract_from_docx(content)
        elif ext in ['.pptx', '.ppt']:
            return TextExtractor.extract_from_pptx(content)
        elif ext in ['.xlsx', '.xls']:
            return TextExtractor.extract_from_xlsx(content)
        elif ext == '.txt':
            return TextExtractor.extract_from_txt(content)
        else:
            raise ValueError(f"サポートされていない拡張子: {ext}")


class TextChunker:
    """テキストをチャンクに分割"""
    
    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200):
        """
        Args:
            chunk_size: チャンクの文字数
            chunk_overlap: チャンク間の重複文字数
        """
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
    
    def split_text(self, text: str, metadata: Optional[Dict] = None) -> List[TextChunk]:
        """
        テキストをチャンクに分割
        
        Args:
            text: 分割するテキスト
            metadata: チャンクに付加するメタデータ
        
        Returns:
            TextChunk のリスト
        """
        if not text or not text.strip():
            return []
        
        # 段落単位で分割を試みる
        paragraphs = self._split_into_paragraphs(text)
        
        chunks = []
        current_chunk = ""
        current_start = 0
        chunk_index = 0
        
        for para in paragraphs:
            # 段落が長すぎる場合は文単位で分割
            if len(para) > self.chunk_size:
                sentences = self._split_into_sentences(para)
                for sentence in sentences:
                    if len(current_chunk) + len(sentence) <= self.chunk_size:
                        current_chunk += sentence + " "
                    else:
                        if current_chunk:
                            chunks.append(self._create_chunk(
                                current_chunk.strip(),
                                chunk_index,
                                current_start,
                                current_start + len(current_chunk),
                                metadata
                            ))
                            chunk_index += 1
                            
                            # オーバーラップ部分を次のチャンクの開始に
                            overlap_text = current_chunk[-self.chunk_overlap:] if len(current_chunk) > self.chunk_overlap else current_chunk
                            current_chunk = overlap_text + sentence + " "
                            current_start += len(current_chunk) - len(overlap_text)
                        else:
                            current_chunk = sentence + " "
            else:
                # 段落全体を追加できるか確認
                if len(current_chunk) + len(para) <= self.chunk_size:
                    current_chunk += para + "\n\n"
                else:
                    # 現在のチャンクを保存
                    if current_chunk:
                        chunks.append(self._create_chunk(
                            current_chunk.strip(),
                            chunk_index,
                            current_start,
                            current_start + len(current_chunk),
                            metadata
                        ))
                        chunk_index += 1
                        
                        # オーバーラップ
                        overlap_text = current_chunk[-self.chunk_overlap:] if len(current_chunk) > self.chunk_overlap else current_chunk
                        current_chunk = overlap_text + para + "\n\n"
                        current_start += len(current_chunk) - len(overlap_text)
                    else:
                        current_chunk = para + "\n\n"
        
        # 最後のチャンク
        if current_chunk.strip():
            chunks.append(self._create_chunk(
                current_chunk.strip(),
                chunk_index,
                current_start,
                current_start + len(current_chunk),
                metadata
            ))
        
        return chunks
    
    def _create_chunk(self, text: str, index: int, start: int, end: int, metadata: Optional[Dict]) -> TextChunk:
        """TextChunk オブジェクトを作成"""
        return TextChunk(
            text=text,
            chunk_index=index,
            start_char=start,
            end_char=end,
            metadata=metadata or {}
        )
    
    @staticmethod
    def _split_into_paragraphs(text: str) -> List[str]:
        """テキストを段落に分割"""
        # 2つ以上の改行で分割
        paragraphs = re.split(r'\n\s*\n', text)
        return [p.strip() for p in paragraphs if p.strip()]
    
    @staticmethod
    def _split_into_sentences(text: str) -> List[str]:
        """テキストを文に分割"""
        # 簡易的な文分割 (日本語と英語)
        sentences = re.split(r'[。.!?]\s*', text)
        return [s.strip() + '。' for s in sentences if s.strip()]


def create_document_id(site_id: str, drive_id: str, item_id: str, chunk_index: Optional[int] = None) -> str:
    """
    ドキュメントIDを生成 (Azure AI Search の制限に準拠)
    
    Args:
        site_id: サイトID
        drive_id: ドライブID
        item_id: アイテムID
        chunk_index: チャンクインデックス (チャンク化する場合)
    
    Returns:
        一意のドキュメントID (英数字、ハイフン、アンダースコアのみ、1024文字以下)
    """
    def sanitize_id_part(text: str, max_length: int = 30) -> str:
        """IDの一部を安全な文字列に変換"""
        # 英数字とハイフン、アンダースコア以外を除去
        safe_text = re.sub(r'[^a-zA-Z0-9\-_]', '', text)
        
        # 先頭が数字や記号の場合は 'doc' を付加
        if safe_text and not safe_text[0].isalpha():
            safe_text = 'doc' + safe_text
        
        # 空の場合はデフォルト値
        if not safe_text:
            safe_text = 'unknown'
        
        # 長さ制限
        if len(safe_text) > max_length:
            safe_text = safe_text[:max_length]
        
        # 末尾のハイフンやアンダースコアを除去
        safe_text = safe_text.rstrip('-_')
        
        # まだ空の場合
        if not safe_text:
            safe_text = 'unknown'
        
        return safe_text
    
    # 各部分を安全な文字列に変換
    safe_site = sanitize_id_part(site_id, 20)
    safe_drive = sanitize_id_part(drive_id, 20) 
    safe_item = sanitize_id_part(item_id, 40)
    
    # IDを構築
    if chunk_index is not None:
        doc_id = f"{safe_site}_{safe_drive}_{safe_item}_{chunk_index}"
    else:
        doc_id = f"{safe_site}_{safe_drive}_{safe_item}"
    
    # 最大長チェック (Azure AI Search は1024文字まで)
    if len(doc_id) > 1024:
        # ハッシュを使って短縮
        import hashlib
        hash_suffix = hashlib.md5(doc_id.encode()).hexdigest()[:8]
        if chunk_index is not None:
            doc_id = f"{safe_site[:10]}_{safe_drive[:10]}_{hash_suffix}_{chunk_index}"
        else:
            doc_id = f"{safe_site[:10]}_{safe_drive[:10]}_{hash_suffix}"
    
    return doc_id


def extract_acl_from_permissions(permissions: List[Dict]) -> Tuple[List[str], List[str]]:
    """
    Graph API の permissions から ACL を抽出
    
    Args:
        permissions: Graph API の permissions レスポンス
    
    Returns:
        (user_ids, group_ids) のタプル
    """
    user_ids = []
    group_ids = []
    
    for perm in permissions:
        # grantedToV2 または grantedTo から情報を取得
        granted = perm.get("grantedToV2") or perm.get("grantedTo")
        
        if granted:
            # ユーザー
            user = granted.get("user")
            if user and "id" in user:
                user_ids.append(user["id"])
            
            # グループ
            group = granted.get("group")
            if group and "id" in group:
                group_ids.append(group["id"])
        
        # grantedToIdentitiesV2 (複数の場合)
        identities = perm.get("grantedToIdentitiesV2") or perm.get("grantedToIdentities", [])
        for identity in identities:
            user = identity.get("user")
            if user and "id" in user:
                user_ids.append(user["id"])
            
            group = identity.get("group")
            if group and "id" in group:
                group_ids.append(group["id"])
    
    # 重複を除去
    return list(set(user_ids)), list(set(group_ids))
